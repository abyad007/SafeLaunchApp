# core/procedure_reader.py
# ─────────────────────────────────────────────────────────────────────────────
# Auto-detects any .xlsx / .pptx file dropped in data/procedures/
# and extracts only quality action items from it.
#
# Design principles (TISAX-compliant):
#   - Zero external calls — all processing is local
#   - Zero new dependencies — uses only python-pptx and openpyxl (already installed)
#   - Fully auditable rule-based extraction — no black-box ML
#   - Preview before use — user validates extracted items before they enter checklist
#
# Extraction strategy:
#   PPTX: reads TABLE cells first (where procedure steps live in Versigent docs),
#         then text frames (for supplementary steps). Filters by linguistic rules.
#   XLSX: detects which column holds action text (longest meaningful strings),
#         reads row by row, applies same linguistic filter.
# ─────────────────────────────────────────────────────────────────────────────

from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Optional, Tuple
import re

# ── Paths ─────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).parent.parent
DATA = ROOT / "data" / "procedures"

# ── Supported extensions ──────────────────────────────────────────────────────
XLSX_EXTS = {".xlsx", ".xlsm", ".xls"}
PPTX_EXTS = {".pptx", ".pptm", ".ppt"}
ALL_EXTS  = XLSX_EXTS | PPTX_EXTS

# ── Known files already mapped to programs (skip from "unassigned") ───────────
KNOWN_FILES = {
    "EAGP_4-4_MG_01-F01_EN.xlsx",
    "EAGP_4-4_MG_01_EN.pptx",
    "EAGP_4-4_CS_01_EN__1_.pptm",
    "EAGP_2-6_BU_02-F01_EN.xlsm",
    "EANP_4-1_CS_01-03_EN.pptx",
    "EAEP_4-1_ME-EDS_10-01_EN.pptx",
    "HOGP_5-1_MG-EDS_01-F01_EN.xlsx",
    "EAGP_5-3_ME_02_EN.pptx",
}

# ─────────────────────────────────────────────────────────────────────────────
# DATA MODEL
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class ExtractedItem:
    """One item extracted from a procedure file."""
    text:     str
    phase:    str  = ""       # slide title or sheet name
    step:     str  = ""       # step number if detected
    critical: bool = False
    source:   str  = ""       # filename it came from


@dataclass
class ExtractionResult:
    """Result of reading one procedure file."""
    filename:      str
    filepath:      Path
    file_type:     str            # "pptx" or "xlsx"
    items:         List[ExtractedItem] = field(default_factory=list)
    error:         Optional[str] = None
    source_method: str = "rule-based"   # "llm" or "rule-based"

    @property
    def item_count(self) -> int:
        return len(self.items)

    @property
    def phases(self) -> List[str]:
        seen, out = set(), []
        for it in self.items:
            if it.phase and it.phase not in seen:
                out.append(it.phase); seen.add(it.phase)
        return out


# ─────────────────────────────────────────────────────────────────────────────
# LINGUISTIC FILTER
# Pure rule-based — no ML, no external calls, fully auditable
# ─────────────────────────────────────────────────────────────────────────────

# Quality action verbs — first word of sentence must be one of these
# OR sentence must contain one of CONTENT_KEYWORDS
ACTION_VERBS = {
    "define","confirm","validate","verify","perform","complete","submit",
    "prepare","obtain","identify","update","notify","communicate","plan",
    "assess","approve","train","release","audit","monitor","implement",
    "develop","document","distribute","inform","request","coordinate",
    "establish","execute","review","check","ensure","gather","receive",
    "analyze","calculate","consolidate","follow","provide","send","set",
    "sign","start","test","transfer","use","conduct","create","support",
    "align","build","close","collect","deploy","detect","enforce","evaluate",
    "initiate","launch","maintain","manage","measure","negotiate","perform",
    "present","process","produce","record","report","resolve","schedule",
    "secure","select","track","update","validate","communicate","escalate",
}

# Content words — if present, likely a quality action even without verb-first
QUALITY_KEYWORDS = {
    "pfmea","control plan","ppap","capability","msa","gauge","cpk","cp",
    "containment","firewall","lpa","audit","ramp-up","run @ rate","pra",
    "special characteristic","kpc","kcc","eol","safe launch","release",
    "approval","inspection","validation","certification","training",
    "process control","quality","defect","non-conformance","escalation",
    "corrective action","root cause","8d","fmea","spc","pre-control",
}

# Patterns that guarantee rejection — not actionable steps
REJECT_PATTERNS = [
    r"^\s*$",                           # empty
    r"^[A-Z\s/&]{1,35}$",              # ALL CAPS short = header
    r"^\d+$",                           # pure number
    r"^[A-Z]{1,4}$",                    # acronym cell (R, S, CS, Mfg)
    r"^(yes|no|n/a|ok|done|open)$",    # status cell
    r"^wk\s*\d",                        # week reference Wk 20/26
    r"^\d{1,2}/\d{2,4}$",              # date pattern
    r"^(responsible|accountable|support|inform|consulted|raci)$",
    r"^(input|output|reference|document|record|role|owner|function|status)$",
    r"^(process step|requirements and|what method|what metric)$",
    r"^(legend|note|remarks|comments?|description)[\s:]*$",
    r"^(step|phase|activity|task|action)[\s:]*$",
    r"^\s*(r|s|cs|mfg|pc|hr|ehs|it|qa)\s*$",  # role abbreviations
    r"^if\s+",                          # conditions not actions
    r"^is\s+|^are\s+|^was\s+|^were\s+|^did\s+|^does\s+",  # questions
    r"^this\s+(procedure|applies|document|form)",
    r"^[*]+",                            # footnote markers
    r"^(approval|support|inform)\s*[-–]",  # RACI legend
    r"^[a-zA-Z]\.\s+\w.*[–-]",    # "B. Document name – code"
    r"^[A-Z]{2,6}_\d",                 # starts with code EAGP_4-3...
    r"^%\s+",                           # % Completion of... (metric lines)
    r"^\w{2,4}:\s*XXX",                # placeholder cells "Quality: XXX"
    r"^(safe launch coverage|customer satisfaction|mq\d+ release|full ppap)$",  # KPI labels
    r".*:\s*$",                         # ends with colon only (label)
    r"^\w+\s*[\u2013-]\s*\w.{5,}",         # glossary "QAM – Quality Assurance..."
    r".*XXX.*",                          # placeholder cells with XXX
    r"^quality:\s|^pc&l:\s|^mfg:\s",  # role assignment cells
]

_compiled_rejects = [re.compile(p, re.IGNORECASE) for p in REJECT_PATTERNS]

# Critical item keywords
CRITICAL_KW = {
    "pfmea","ppap","special characteristic","sc","safety","lpa",
    "run @ rate","pra","bmg","containment","100%","firewall","kpc",
    "control plan","msa","cpk","process release","full ppap","eol",
    "critical characteristic","customer approval",
}


def _is_quality_action(text: str) -> bool:
    """
    Returns True if this text looks like a quality action item.
    Pure rule-based — no ML.

    Rules applied in order:
    1. Length check (20–250 chars)
    2. Reject patterns (headers, roles, conditions, status words)
    3. Action verb first word  OR  quality keyword present
    """
    t = text.strip()

    # Rule 1 — length
    if len(t) < 20 or len(t) > 250:
        return False

    tl = t.lower()

    # Rule 2 — reject patterns
    for pattern in _compiled_rejects:
        if pattern.search(tl):
            return False

    # Rule 3a — first word is an action verb
    first = tl.split()[0].rstrip("sed") if tl.split() else ""
    if first in ACTION_VERBS:
        return True

    # Rule 3b — contains quality keyword
    if any(kw in tl for kw in QUALITY_KEYWORDS):
        return True

    # Rule 3c — starts with a digit followed by content (numbered step)
    if re.match(r"^\d+[\s.)]+\w{4,}", t):
        # Only accept if not a pure table-of-contents reference
        if not re.match(r"^\d+[\s.)]+\d", t):
            return True

    return False


def _is_critical(text: str) -> bool:
    tl = text.lower()
    return any(kw in tl for kw in CRITICAL_KW)


def _clean_text(text: str) -> str:
    """Remove leading step numbers, bullets, role prefixes."""
    t = text.strip()
    t = re.sub(r"^\d+[.):\s]+", "", t)           # "1. " or "1) " or "1: "
    t = re.sub(r"(?i)^step\s*\d+[.):]*\s*", "", t)  # "Step 1: "
    t = t.strip("•·–—-●▪ \t").strip()
    return t


def _deduplicate(items: List[ExtractedItem]) -> List[ExtractedItem]:
    """
    Remove near-duplicate items.
    Two items are duplicates if their first 55 chars match (case-insensitive)
    OR they share 5+ content words.
    """
    seen_keys:  set = set()
    seen_words: List[set] = []
    result:     List[ExtractedItem] = []

    for item in items:
        key = item.text[:55].lower().strip()
        if key in seen_keys:
            continue

        # Word-overlap check
        words = set(re.findall(r"\b\w{4,}\b", item.text.lower()))
        is_dup = any(len(words & w) >= 5 for w in seen_words)
        if is_dup:
            continue

        seen_keys.add(key)
        seen_words.append(words)
        result.append(item)

    return result


# ─────────────────────────────────────────────────────────────────────────────
# PPTX READER
# ─────────────────────────────────────────────────────────────────────────────

def _read_pptx(filepath: Path) -> List[ExtractedItem]:
    """
    Extracts quality action items from a PowerPoint file.

    Priority order per slide:
    1. TABLE cells (shape_type=19) — main procedure steps in Versigent docs
    2. Text frames — supplementary steps and decisions

    Slide title (placeholder idx=0) becomes the phase label.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return []

    prs   = Presentation(str(filepath))
    items: List[ExtractedItem] = []
    seen:  set = set()
    step_n = 0

    for slide_idx, slide in enumerate(prs.slides):
        # Get slide title
        phase = ""
        for shape in slide.shapes:
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.idx == 0:
                        phase = shape.text_frame.text.strip()[:60]
                        break
                except Exception:
                    pass

        slide_items: List[str] = []

        # 1) TABLE cells — primary source
        for shape in slide.shapes:
            if shape.shape_type != 19:    # 19 = MSO TABLE
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and text not in seen:
                            slide_items.append(text)
                            seen.add(text)

        # 2) Text frames — secondary source
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Skip the title placeholder — already captured
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.idx == 0:
                        continue
                except Exception:
                    pass
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and text not in seen:
                    slide_items.append(text)
                    seen.add(text)

        # Apply quality filter
        for text in slide_items:
            cleaned = _clean_text(text)
            if _is_quality_action(cleaned):
                step_n += 1
                items.append(ExtractedItem(
                    text=cleaned,
                    phase=phase,
                    step=str(step_n),
                    critical=_is_critical(cleaned),
                    source=filepath.name,
                ))

    return _deduplicate(items)


# ─────────────────────────────────────────────────────────────────────────────
# XLSX READER
# ─────────────────────────────────────────────────────────────────────────────

def _detect_action_column(sheet, max_rows: int = 30) -> int:
    """
    Scans the first max_rows rows to find which column contains
    the longest meaningful text strings — that's the action column.
    Returns 0-based column index, default 0.
    """
    col_scores: dict = {}
    for row in sheet.iter_rows(min_row=1, max_row=max_rows, values_only=True):
        for col_idx, cell_val in enumerate(row):
            if not cell_val:
                continue
            text = str(cell_val).strip()
            # Score = length of meaningful text (penalise very short or ALL-CAPS)
            if len(text) >= 20 and not text.isupper():
                col_scores[col_idx] = col_scores.get(col_idx, 0) + len(text)
    return max(col_scores, key=col_scores.get) if col_scores else 0


def _read_xlsx(filepath: Path, sheet_name: str = "") -> List[ExtractedItem]:
    """
    Extracts quality action items from an Excel workbook.

    If sheet_name is empty, reads all sheets.
    Auto-detects which column contains the action step text.
    """
    try:
        import openpyxl
    except ImportError:
        return []

    wb = openpyxl.load_workbook(str(filepath), read_only=False, data_only=True)
    items: List[ExtractedItem] = []

    sheets_to_read = (
        [sheet_name] if sheet_name and sheet_name in wb.sheetnames
        else [s for s in wb.sheetnames
              if not any(skip in s.lower()
                         for skip in ["template","reference","revision","index","legend"])]
    )

    for sname in sheets_to_read:
        sheet  = wb[sname]
        phase  = sname  # sheet name = phase label
        col_idx = _detect_action_column(sheet)
        step_n  = 0

        for row in sheet.iter_rows(min_row=2, values_only=True):
            if not row or col_idx >= len(row):
                continue
            cell_val = row[col_idx]
            if not cell_val:
                continue

            # Handle multi-line cell values
            lines = str(cell_val).splitlines()
            for line in lines:
                cleaned = _clean_text(line)
                if _is_quality_action(cleaned):
                    step_n += 1
                    items.append(ExtractedItem(
                        text=cleaned,
                        phase=phase,
                        step=f"{sname[:3].upper()}-{step_n}",
                        critical=_is_critical(cleaned),
                        source=filepath.name,
                    ))

    wb.close()
    return _deduplicate(items)


# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API
# ─────────────────────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────────────────────────────────────
# AUTO-CLASSIFICATION
# Determines which program type a file belongs to without user input.
# Uses filename patterns first (fast), then falls back to content keywords.
# ─────────────────────────────────────────────────────────────────────────────

# Filename fragments → program type (checked in order, first match wins)
_FILENAME_RULES: List[Tuple[str, str]] = [
    # Business Transfer
    ("eagp_4-4",    "transfer"),
    ("eagp_2-6",    "transfer"),
    ("bu_02",       "transfer"),
    ("transfer",    "transfer"),
    ("business",    "transfer"),
    # Restart After Shutdown
    ("hogp_5-1",    "restart"),
    ("shutdown",    "restart"),
    ("restart",     "restart"),
    # Engineering / MY Change
    ("eaep_4-1",    "my_change"),
    ("engineering", "my_change"),
    ("my_change",   "my_change"),
    ("eccl",        "my_change"),
    # New Program
    ("eanp_4-1",    "new"),
    ("new_prog",    "new"),
    ("apqp",        "new"),
    ("launch",      "new"),
    # Capacity
    ("eagp_5-3",    "capacity"),
    ("capacity",    "capacity"),
    ("ramp",        "capacity"),
    # Absence / Turnover
    ("absence",     "absence"),
    ("turnover",    "absence"),
    ("workforce",   "absence"),
]

# Content keywords → program type (used when filename gives no match)
# Each entry: (keyword, program, weight)
_CONTENT_RULES: List[Tuple[str, str, int]] = [
    # Business Transfer signals
    ("business transfer",        "transfer", 5),
    ("receiving plant",          "transfer", 4),
    ("transfer checklist",       "transfer", 4),
    ("pra",                      "transfer", 3),
    ("safety stock",             "transfer", 2),
    ("31 steps",                 "transfer", 2),
    # Restart signals
    ("shutdown prep",            "restart",  5),
    ("restart",                  "restart",  4),
    ("shutdown",                 "restart",  3),
    ("lpa.*area",                "restart",  3),
    ("incoming area",            "restart",  3),
    # Engineering Change signals
    ("eccl",                     "my_change", 5),
    ("engineering change",       "my_change", 4),
    ("model year",               "my_change", 4),
    ("change notice",            "my_change", 3),
    ("e&o",                      "my_change", 3),
    # New Program signals
    ("mq1",                      "new",      5),
    ("mq2",                      "new",      5),
    ("new program",              "new",      4),
    ("start of production",      "new",      3),
    ("sop",                      "new",      2),
    # Capacity signals
    ("capacity change",          "capacity", 5),
    ("atv",                      "capacity", 4),
    ("ramp.up curve",            "capacity", 3),
    ("volume change",            "capacity", 3),
    # Absence signals
    ("turnover",                 "absence",  5),
    ("absenteeism",              "absence",  4),
    ("headcount gap",            "absence",  4),
]

_PROG_LABELS = {
    "new":       "New Program",
    "transfer":  "Business Transfer",
    "my_change": "MY / Engineering Change",
    "restart":   "Restart After Shutdown",
    "absence":   "Absence / Turnover",
    "capacity":  "Capacity Change",
}


def auto_classify_file(filepath: Path) -> Tuple[str, float]:
    """
    Determines the program type for a procedure file automatically.

    Returns (program_type, confidence) where confidence is 0.0–1.0.
    Returns ("unknown", 0.0) if no match found.

    Strategy:
    1. Check filename against _FILENAME_RULES (fast, usually enough)
    2. If no filename match, sample file content and score against _CONTENT_RULES
    """
    fname = filepath.name.lower()

    # Step 1 — filename rules (O(n) scan, n ≤ 25)
    for fragment, prog in _FILENAME_RULES:
        if fragment in fname:
            return prog, 1.0

    # Step 2 — content scan (only if filename gives no signal)
    text_sample = _sample_file_text(filepath, max_chars=4000)
    if not text_sample:
        return "unknown", 0.0

    tl = text_sample.lower()
    scores: dict = {}   # prog → cumulative weight

    for keyword, prog, weight in _CONTENT_RULES:
        import re as _re
        if _re.search(keyword, tl):
            scores[prog] = scores.get(prog, 0) + weight

    if not scores:
        return "unknown", 0.0

    best_prog  = max(scores, key=scores.get)
    best_score = scores[best_prog]
    total      = sum(scores.values())
    confidence = min(best_score / max(total, 1), 1.0)

    # Only return a classification if we have meaningful confidence
    if best_score < 4:
        return "unknown", round(confidence, 2)

    return best_prog, round(confidence, 2)


def _sample_file_text(filepath: Path, max_chars: int = 4000) -> str:
    """
    Reads a small sample of text from a file for classification.
    Stops after max_chars to keep it fast.
    """
    ext = filepath.suffix.lower()
    text_parts: List[str] = []
    chars_read  = 0

    try:
        if ext in PPTX_EXTS:
            from pptx import Presentation
            prs = Presentation(str(filepath))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            text_parts.append(t)
                            chars_read += len(t)
                            if chars_read >= max_chars:
                                break
                if chars_read >= max_chars:
                    break

        elif ext in XLSX_EXTS:
            import openpyxl
            wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
            for sname in wb.sheetnames[:5]:
                sheet = wb[sname]
                for row in sheet.iter_rows(min_row=1, max_row=30, values_only=True):
                    for cell in row:
                        if cell:
                            t = str(cell).strip()
                            if len(t) > 3:
                                text_parts.append(t)
                                chars_read += len(t)
                                if chars_read >= max_chars:
                                    break
                    if chars_read >= max_chars:
                        break
                wb.close()

    except Exception:
        pass

    return " ".join(text_parts)


def scan_new_files() -> List[Path]:
    """
    Scans data/procedures/ for any file not in the KNOWN_FILES set.
    Returns a list of Path objects for unrecognised procedure files.
    """
    if not DATA.exists():
        return []
    return [
        f for f in DATA.iterdir()
        if f.suffix.lower() in ALL_EXTS
        and f.name not in KNOWN_FILES
    ]


def extract_from_file(filepath: Path) -> ExtractionResult:
    """
    Main entry point. Pass any .xlsx or .pptx Path.

    Extraction strategy (automatic, no configuration needed at call site):
      1. Try LLM extraction if LLM server is configured and reachable
         → Uses vLLM or Ollama at the endpoint in config.py
         → Returns clean, structured items via the language model
      2. Fall back to rule-based extraction if LLM is unavailable or fails
         → Uses python-pptx / openpyxl with linguistic filters
         → Zero dependencies on external services

    The caller never needs to know which method was used —
    ExtractionResult.source_method reports it for display.

    Usage:
        result = extract_from_file(Path("data/procedures/MY_PROC.pptx"))
        print(result.source_method)   # "llm" or "rule-based"
        for item in result.items:
            print(item.step, item.text)
    """
    ext = filepath.suffix.lower()

    if not filepath.exists():
        return ExtractionResult(
            filename=filepath.name,
            filepath=filepath,
            file_type="unknown",
            error=f"File not found: {filepath}",
        )

    file_type = (
        "pptx" if ext in PPTX_EXTS else
        "xlsx" if ext in XLSX_EXTS else
        "unknown"
    )

    if file_type == "unknown":
        return ExtractionResult(
            filename=filepath.name,
            filepath=filepath,
            file_type="unknown",
            error=f"Unsupported file type: {ext}",
        )

    # ── Step 1: Try LLM extraction ────────────────────────────────────────
    try:
        from core.llm_extractor import extract_with_llm
        llm_raw = extract_with_llm(filepath)

        if llm_raw is not None:
            # Convert LLMItems to ExtractedItems
            items = [
                ExtractedItem(
                    text=it.text,
                    phase=it.phase,
                    step=it.step,
                    critical=it.critical,
                    source=filepath.name,
                )
                for it in llm_raw
            ]
            result = ExtractionResult(
                filename=filepath.name,
                filepath=filepath,
                file_type=file_type,
                items=items,
            )
            result.source_method = "llm"
            return result

    except Exception as e:
        # LLM attempt failed for any reason — continue to rule-based
        print(f"[procedure_reader] LLM extraction error: {e} — falling back")

    # ── Step 2: Rule-based extraction (fallback) ──────────────────────────
    try:
        if file_type == "pptx":
            items = _read_pptx(filepath)
        else:
            items = _read_xlsx(filepath)

        result = ExtractionResult(
            filename=filepath.name,
            filepath=filepath,
            file_type=file_type,
            items=items,
        )
        result.source_method = "rule-based"
        return result

    except Exception as e:
        return ExtractionResult(
            filename=filepath.name,
            filepath=filepath,
            file_type=file_type,
            error=str(e),
        )


def extraction_summary(result: ExtractionResult) -> str:
    """Returns a one-line summary for the sidebar status panel."""
    if result.error:
        return f"Error: {result.error[:60]}"
    return (
        f"{result.item_count} quality items extracted "
        f"across {len(result.phases)} phase(s)"
    )
