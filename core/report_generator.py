# core/report_generator.py
# ─────────────────────────────────────────────────────────────────────────────
# Generates PPT, PDF (via HTML→PDF), and Excel reports from plan data.
#
# Libraries used:
#   python-pptx  : builds PowerPoint files programmatically
#   openpyxl     : builds Excel files
#   reportlab    : builds PDF from scratch
#
# Key python-pptx concepts:
#   Presentation()            : loads or creates a PPTX
#   prs.slides.add_slide()    : adds a new slide
#   slide.shapes.add_textbox(): places a text box at exact position
#   Inches(n)                 : converts inches to EMU (PPT unit)
#   Pt(n)                     : converts points to EMU for font sizes
#   RGBColor(r, g, b)         : color from 0-255 values
# ─────────────────────────────────────────────────────────────────────────────

from pathlib import Path
from datetime import datetime
from typing import List
import io           # io.BytesIO = in-memory binary stream (no temp files needed)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from core.scoring_engine import ScoreResult
from core.checklist_loader import ChecklistItem, get_phases
from core.theme import token as _tok   # design-system tokens (single source of truth)

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils  import get_column_letter

# ── Paths ─────────────────────────────────────────────────────────────────────
ROOT     = Path(__file__).parent.parent
TEMPLATE   = ROOT / "data" / "templates" / "Versigent_MasterTemplate.pptx"
LOGO_COLOR = ROOT / "data" / "templates" / "versigent_logo.png"       # copper on transparent
LOGO_WHITE = ROOT / "data" / "templates" / "versigent_logo_white.png" # white for dark bg
OUTPUT     = ROOT / "output"

# ── Brand palette & fonts — derived from the design system ───────────────────
# Single source of truth: design/tokens.json (via core.theme). Same constant
# names as before, so the rest of this module is untouched — only the values now
# come from the tokens, so PPT/Excel exports match the app exactly.

def _trgb(path: str) -> RGBColor:
    """Token color path -> RGBColor (python-pptx)."""
    h = _tok(path).lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _thex(path: str) -> str:
    """Token color path -> bare 'RRGGBB' hex (openpyxl)."""
    return _tok(path).lstrip("#").upper()

# Typography (tokens: display = Barlow Condensed, body = Barlow)
DISPLAY_FONT = "Barlow Condensed"
BODY_FONT    = "Barlow"

# Brand
NAVY      = _trgb("primitive.color.navy.900")
COPPER    = _trgb("primitive.color.copper.500")
GOLD      = _trgb("primitive.color.gold.400")
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
# Neutrals
DARK_GRAY = _trgb("primitive.color.gray.900")
GRAY      = _trgb("semantic.color.muted-foreground")
LIGHT_GRAY= _trgb("primitive.color.stone.200")
# Status — the -text (700) tier reads AA as text on white AND keeps white text
# legible when used as a fill (badges/score block).
GREEN     = _trgb("semantic.color.risk-low-text")
RED       = _trgb("semantic.color.risk-high-text")

RISK_COLORS = {
    "HIGH":   _trgb("semantic.color.risk-high-text"),
    "MEDIUM": _trgb("semantic.color.risk-medium-text"),
    "LOW":    _trgb("semantic.color.risk-low-text"),
}

# Hex strings (no '#') for openpyxl, from the same tokens
HX_NAVY   = _thex("primitive.color.navy.900")
HX_COPPER = _thex("primitive.color.copper.500")
HX_FG     = _thex("primitive.color.gray.900")
HX_MUTED  = _thex("semantic.color.muted-foreground")
HX_BORDER = _thex("primitive.color.stone.200")
HX_HIGH_T = _thex("semantic.color.risk-high-text")
HX_MED_T  = _thex("semantic.color.risk-medium-text")
HX_LOW_T  = _thex("semantic.color.risk-low-text")

# Program display names
PROG_LABELS = {
    "new":       "New Program",
    "new_plant": "New Plant Launch",
    "transfer":  "Business Transfer",
    "my_change": "Engineering Change",
    "restart":   "Restart After Shutdown",
    "absence":   "Absence / Turnover",
    "capacity":  "Capacity Change",
}

PROG_TITLES = {
    "new":       "NEW PROGRAM SAFE LAUNCH",
    "new_plant": "NEW PLANT FLAWLESS LAUNCH",
    "transfer":  "BUSINESS TRANSFER PLAN",
    "my_change": "ENGINEERING CHANGE PLAN",
    "restart":   "RESTART AFTER SHUTDOWN",
    "absence":   "ABSENCE / TURNOVER PLAN",
    "capacity":  "CAPACITY CHANGE PLAN",
}

PROG_PROCS = {
    "new":       "EANP_4-1_CS_01-03_EN",
    "new_plant": "Aptiv New Plant Flawless Launch Initiative | RASIC Apr 2025",
    "transfer":  "EAGP_4-4_MG_01 | EAGP_2-6_BU_02",
    "my_change": "EAEP_4-1_ME-EDS_10-01_EN",
    "restart":   "HOGP_5-1_MG-EDS_01-F01_EN",
    "absence":   "Workforce Stability Program",
    "capacity":  "EAGP_5-3_ME_02_EN",
}


# ═════════════════════════════════════════════════════════════════════════════
# INTERNAL HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def _add_text(slide, text: str, left, top, width, height,
              font_size=12, bold=False, color=DARK_GRAY,
              align=PP_ALIGN.LEFT, font_name=BODY_FONT):
    """
    Adds a text box to a slide. All position/size args should be Inches() values.
    Returns the text frame so you can add more runs if needed.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True             # wrap long text automatically

    para = tf.paragraphs[0]
    para.alignment = align

    run = para.add_run()            # a "run" is a text segment with formatting
    run.text = str(text)
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.name   = font_name

    return tf


def _add_rect(slide, left, top, width, height, fill_rgb: RGBColor, line_rgb: RGBColor = None):
    """
    Adds a filled rectangle shape to a slide.
    MSO_SHAPE_TYPE.RECTANGLE is a built-in constant.
    """
    shape = slide.shapes.add_shape(
        1,      # 1 = MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()                     # solid fill (not gradient/pattern)
    shape.fill.fore_color.rgb = fill_rgb

    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()       # no border

    return shape


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Converts a hex string like 'CD7925' to RGBColor(0xCD, 0x79, 0x25)."""
    h = hex_str.lstrip("#")   # remove # if present
    # int(s, 16) converts hex string to int
    # [0:2] slices first 2 chars, [2:4] middle 2, [4:6] last 2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _footer(slide, page_num: int, part: str, prog_type: str):
    """Adds consistent footer to every slide."""
    year = datetime.now().year
    _add_text(
        slide, f"Safe Launch Plan | {year} | Versigent Confidential",
        Inches(0.4), Inches(7.2), Inches(9), Inches(0.25),
        font_size=8, color=GRAY,
    )
    _add_text(
        slide, str(page_num),
        Inches(12.8), Inches(7.2), Inches(0.4), Inches(0.25),
        font_size=8, color=GRAY, align=PP_ALIGN.RIGHT,
    )


# ═════════════════════════════════════════════════════════════════════════════
# PPT GENERATOR
# ═════════════════════════════════════════════════════════════════════════════

def generate_ppt(
    prog_type:  str,
    part_name:  str,
    result:     ScoreResult,
    checklist:  List[ChecklistItem],
    meta:       dict,              # extra data like dates, volumes, headcount
    customer_name: str = "Other",
) -> bytes:
    """
    Builds a complete PowerPoint and returns it as bytes.

    Returning bytes (instead of saving to file) lets Streamlit offer
    an in-memory download — no temp files needed.

    Parameters:
        prog_type    : "new" | "transfer" | "restart" | ...
        part_name    : part/project name entered by user
        result       : ScoreResult from scoring_engine
        checklist    : List[ChecklistItem] from checklist_loader
        meta         : dict with dates, volumes, headcount etc.
        customer_name: display name of the OEM
    """

    # ── Load template — use it directly so all branding is inherited ─────
    # Strategy: open the Versigent template file, add our slides to it
    # using its "Blank Layout", then delete the original template slides.
    # This preserves fonts, colors, wave art, logo, and theme perfectly.

    if TEMPLATE.exists():
        prs = Presentation(str(TEMPLATE))
        # Remember how many slides the template has — we delete these at the end
        _original_slide_count = len(prs.slides)
        # Find the Blank Layout (index 5 in Versigent template)
        blank_layout = None
        for _lay in prs.slide_layouts:
            if "blank" in _lay.name.lower():
                blank_layout = _lay
                break
        if blank_layout is None:
            blank_layout = prs.slide_layouts[5]
    else:
        # No template found — create from scratch with correct dimensions
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        _original_slide_count = 0
        blank_layout = prs.slide_layouts[0]

    def _add_branded_slide():
        """Adds a slide using the template's Blank Layout."""
        return prs.slides.add_slide(blank_layout)

    prog_label = PROG_LABELS.get(prog_type, "Safe Launch")
    prog_title = PROG_TITLES.get(prog_type, "SAFE LAUNCH PLAN")
    prog_proc  = PROG_PROCS.get(prog_type, "")
    risk_color = RISK_COLORS.get(result.risk, GRAY)

    # ── SLIDE 1: Cover ────────────────────────────────────────────────────────
    s1 = _add_branded_slide()

    # Navy background
    _add_rect(s1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), NAVY, NAVY)

    # Copper accent bar
    _add_rect(s1, Inches(0), Inches(0), Inches(0.08), Inches(7.5), COPPER, COPPER)

    # Program title
    _add_text(s1, prog_title,
              Inches(0.5), Inches(1.2), Inches(9), Inches(0.9),
              font_size=38, bold=True, color=WHITE)

    # Part name in gold
    _add_text(s1, part_name,
              Inches(0.5), Inches(2.05), Inches(9), Inches(0.6),
              font_size=22, bold=False, color=GOLD)

    # Risk badge — colored rectangle with text
    _add_rect(s1, Inches(0.5), Inches(2.9), Inches(2.2), Inches(0.45),
              risk_color, risk_color)
    _add_text(s1, f"{result.risk} RISK — {result.score}/100",
              Inches(0.5), Inches(2.92), Inches(2.2), Inches(0.41),
              font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Procedure reference
    _add_text(s1, f"Procedure: {prog_proc}",
              Inches(0.5), Inches(3.5), Inches(9), Inches(0.35),
              font_size=11, color=_hex_to_rgb("8BAABF"))

    # Metrics tiles (6 cells in 2 rows of 3)
    kpis = _build_kpis(prog_type, result, meta, customer_name)
    for i, kpi in enumerate(kpis[:6]):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.1)
        y = Inches(4.1 + row * 0.9)
        _add_rect(s1, x, y, Inches(3.85), Inches(0.78),
                  _hex_to_rgb("1E3252"), COPPER)
        _add_text(s1, kpi["label"],
                  x + Inches(0.1), y + Inches(0.04), Inches(3.65), Inches(0.22),
                  font_size=7, color=_hex_to_rgb("CB9552"))
        _add_text(s1, kpi["value"],
                  x + Inches(0.1), y + Inches(0.26), Inches(3.65), Inches(0.4),
                  font_size=16, bold=True, color=WHITE)

    _footer(s1, 1, part_name, prog_type)

    # ── SLIDE 2: Risk Analysis ────────────────────────────────────────────────
    s2 = _add_branded_slide()
    _add_rect(s2, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
              _hex_to_rgb("FAFAF8"), _hex_to_rgb("FAFAF8"))
    _add_rect(s2, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.03), COPPER, COPPER)

    _add_text(s2, f"{part_name} — {prog_label.upper()} RISK ANALYSIS",
              Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.55),
              font_size=22, bold=True, color=NAVY)

    # Score block
    _add_rect(s2, Inches(0.4), Inches(1.0), Inches(2.5), Inches(2.2), risk_color, risk_color)
    _add_text(s2, str(result.score),
              Inches(0.4), Inches(1.1), Inches(2.5), Inches(1.4),
              font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s2, "RISK SCORE / 100",
              Inches(0.4), Inches(2.5), Inches(2.5), Inches(0.35),
              font_size=8, color=_hex_to_rgb("FFD080"), align=PP_ALIGN.CENTER)
    _add_text(s2, f"{result.risk} RISK",
              Inches(0.4), Inches(2.82), Inches(2.5), Inches(0.38),
              font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Recommendation box
    _add_rect(s2, Inches(3.1), Inches(1.0), Inches(9.8), Inches(2.2),
              _hex_to_rgb("F8F5F0"), LIGHT_GRAY)
    _add_text(s2, "RECOMMENDATION",
              Inches(3.3), Inches(1.1), Inches(9.4), Inches(0.3),
              font_size=9, bold=True, color=COPPER)
    _add_text(s2, result.recommendation,
              Inches(3.3), Inches(1.45), Inches(9.4), Inches(1.6),
              font_size=11, color=DARK_GRAY)

    # 6 KPI tiles
    kpis = _build_kpis(prog_type, result, meta, customer_name)
    for i, kpi in enumerate(kpis[:6]):
        col = i % 3
        row = i // 3
        x = Inches(0.4 + col * 4.1)
        y = Inches(3.35 + row * 1.15)
        _add_rect(s2, x, y, Inches(3.85), Inches(1.0), WHITE, LIGHT_GRAY)
        _add_rect(s2, x, y, Inches(3.85), Inches(0.05), COPPER, COPPER)
        _add_text(s2, kpi["label"], x+Inches(0.1), y+Inches(0.08), Inches(3.65), Inches(0.22), font_size=7, color=GRAY)
        _add_text(s2, kpi["value"], x+Inches(0.1), y+Inches(0.3),  Inches(3.65), Inches(0.42), font_size=18, bold=True, color=NAVY)
        _add_text(s2, kpi["sub"],   x+Inches(0.1), y+Inches(0.72), Inches(3.65), Inches(0.22), font_size=8, color=GRAY)

    _footer(s2, 2, part_name, prog_type)

    # ── SLIDE 3: Risk Factors ─────────────────────────────────────────────────
    s3 = _add_branded_slide()
    _add_rect(s3, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
              _hex_to_rgb("FAFAF8"), _hex_to_rgb("FAFAF8"))
    _add_text(s3, f"{part_name} — RISK FACTOR BREAKDOWN",
              Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.55),
              font_size=22, bold=True, color=NAVY)
    _add_rect(s3, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.03), COPPER, COPPER)

    # Bar chart of each factor
    bar_colors = [COPPER, _hex_to_rgb("E09030"), _hex_to_rgb("A86618"),
                  _hex_to_rgb("D49030"), _hex_to_rgb("8B5E10"), GOLD]
    for i, factor in enumerate(result.factors):
        y = Inches(1.1 + i * 0.45)
        if y > Inches(7.0):
            break
        pct = factor.value / factor.max if factor.max else 0
        # Factor name label
        _add_text(s3, f"{factor.name}",
                  Inches(0.4), y, Inches(3.8), Inches(0.35),
                  font_size=10, color=DARK_GRAY)
        # Score value
        _add_text(s3, f"{factor.value}/{factor.max}",
                  Inches(4.3), y, Inches(0.7), Inches(0.35),
                  font_size=10, bold=True, color=NAVY)
        # Background bar
        _add_rect(s3, Inches(5.1), y + Inches(0.08),
                  Inches(7.5), Inches(0.2),
                  _hex_to_rgb("E5DFD3"), _hex_to_rgb("E5DFD3"))
        # Fill bar — width proportional to score
        if pct > 0:
            _add_rect(s3, Inches(5.1), y + Inches(0.08),
                      Inches(7.5 * pct), Inches(0.2),
                      bar_colors[i % len(bar_colors)],
                      bar_colors[i % len(bar_colors)])
        # Percent label
        _add_text(s3, f"{factor.percent}%",
                  Inches(12.7), y, Inches(0.5), Inches(0.35),
                  font_size=9, color=GRAY)

    # Versigent logo on slide 3
    if LOGO_COLOR.exists():
        s3.shapes.add_picture(str(LOGO_COLOR), Inches(10.3), Inches(6.98), Inches(2.6), Inches(0.53))

    _footer(s3, 3, part_name, prog_type)

    # ── SLIDE 4: Quality Domain Overview — 4-column split ────────────────────
    # One column per quality phase — mirrors the MO2 Quality Performance Review
    s4 = _add_branded_slide()
    _add_rect(s4, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
              _hex_to_rgb("FAFAF8"), _hex_to_rgb("FAFAF8"))
    _add_text(s4, f"{part_name} — {prog_label.upper()} QUALITY PLAN",
              Inches(0.3), Inches(0.18), Inches(12.7), Inches(0.52),
              font_size=22, bold=True, color=NAVY)
    _add_rect(s4, Inches(0.3), Inches(0.72), Inches(12.73), Inches(0.03), COPPER, COPPER)

    PHASE_COLORS_4 = {
        "Quality Structure":     "B91C1C",
        "Customer Requirements": "1565A0",
        "Workforce Actions":     "B45309",
        "Shutdown Prep":         "16283F",
        "Quality Readiness":     "16283F",
        "Quality Confirmation":  "1A5C1A",
    }
    PHASE_ABBREV = {
        "Quality Structure":     "QS",
        "Customer Requirements": "CR",
        "Workforce Actions":     "WF",
        "Shutdown Prep":         "SP",
        "Quality Readiness":     "QR",
        "Quality Confirmation":  "QC",
    }

    phases4   = get_phases(checklist)
    n_cols    = len(phases4)
    col_w     = 12.73 / n_cols
    start_x   = 0.3
    header_y  = 0.78
    header_h  = 0.52
    content_y = header_y + header_h
    content_h = 5.8
    row_h     = 0.34

    for col_idx, phase in enumerate(phases4):
        cx        = Inches(start_x + col_idx * col_w)
        cw        = Inches(col_w - 0.06)
        hdr_color = _hex_to_rgb(PHASE_COLORS_4.get(phase, "16283F"))
        abbrev    = PHASE_ABBREV.get(phase, phase[:2].upper())

        # Column header
        _add_rect(s4, cx, Inches(header_y), cw, Inches(header_h), hdr_color, hdr_color)
        _add_rect(s4, cx + Inches(0.06), Inches(header_y + 0.08),
                  Inches(0.28), Inches(0.28), WHITE, WHITE)
        _add_text(s4, abbrev,
                  cx + Inches(0.06), Inches(header_y + 0.06),
                  Inches(0.28), Inches(0.3),
                  font_size=7, bold=True, color=hdr_color, align=PP_ALIGN.CENTER)
        _add_text(s4, phase.upper(),
                  cx + Inches(0.38), Inches(header_y + 0.08),
                  cw - Inches(0.42), Inches(header_h - 0.12),
                  font_size=8, bold=True, color=WHITE)

        # Column background
        _add_rect(s4, cx, Inches(content_y), cw, Inches(content_h),
                  _hex_to_rgb("FFFFFF"), _hex_to_rgb("E5DFD3"))

        # Items
        phase_items = [it for it in checklist if it.phase == phase]
        for row_idx, item in enumerate(phase_items):
            iy = Inches(content_y + 0.06 + row_idx * row_h)
            if iy + Inches(row_h) > Inches(content_y + content_h - 0.05):
                remaining = len(phase_items) - row_idx
                _add_text(s4, f"+ {remaining} more…",
                          cx + Inches(0.1), iy, cw - Inches(0.12), Inches(0.28),
                          font_size=7, color=GRAY)
                break
            row_bg = _hex_to_rgb("F9F6F2") if row_idx % 2 == 0 else _hex_to_rgb("FFFFFF")
            _add_rect(s4, cx + Inches(0.04), iy,
                      cw - Inches(0.08), Inches(row_h - 0.04),
                      row_bg, _hex_to_rgb("EDE7DA"))
            # Status dot
            if item.done:   dc, ds = GREEN, "✓"
            elif item.critical: dc, ds = RED, "!"
            else:           dc, ds = _hex_to_rgb("64748B"), "○"
            _add_rect(s4, cx + Inches(0.06), iy + Inches(0.07),
                      Inches(0.16), Inches(0.16), dc, dc)
            _add_text(s4, ds,
                      cx + Inches(0.06), iy + Inches(0.04),
                      Inches(0.16), Inches(0.22),
                      font_size=6, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            txt = item.text[:54] + "…" if len(item.text) > 54 else item.text
            _add_text(s4, txt,
                      cx + Inches(0.25), iy + Inches(0.04),
                      cw - Inches(0.30), Inches(row_h - 0.06),
                      font_size=7,
                      color=RED if item.critical else DARK_GRAY)
            # Show owner badge if assigned
            owner_v = getattr(item, "owner", "")
            if owner_v:
                _add_text(s4, f"● {owner_v[:12]}",
                          cx + Inches(0.25), iy + Inches(row_h - 0.1),
                          cw - Inches(0.30), Inches(0.14),
                          font_size=5.5, color=_hex_to_rgb("CD7925"))

        # Item count footer per column
        _add_text(s4, f"{len(phase_items)} items",
                  cx, Inches(content_y + content_h + 0.04),
                  cw, Inches(0.22),
                  font_size=7, color=GRAY, align=PP_ALIGN.CENTER)

    # Legend
    for li, (sym, lbl, lc) in enumerate([("✓","Completed","15803D"),("!","Critical","B91C1C"),("○","Pending","64748B")]):
        lx = Inches(0.3 + li * 2.2)
        ly = Inches(7.04)
        _add_rect(s4, lx, ly, Inches(0.18), Inches(0.18), _hex_to_rgb(lc), _hex_to_rgb(lc))
        _add_text(s4, sym, lx, ly - Inches(0.02), Inches(0.18), Inches(0.22),
                  font_size=6, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(s4, lbl, lx + Inches(0.22), ly, Inches(1.8), Inches(0.2),
                  font_size=8, color=DARK_GRAY)

    # Logo on slide 4
    if LOGO_COLOR.exists():
        s4.shapes.add_picture(str(LOGO_COLOR), Inches(10.3), Inches(7.02), Inches(2.6), Inches(0.53))

    _footer(s4, 4, part_name, prog_type)

    # ── SLIDES 5+: One slide per quality phase, split at 20 items ──────────
    ITEMS_PER_SLIDE = 20
    phases_chk = get_phases(checklist)

    # Build 4-tuples: (items, phase_name, page_within_phase, total_pages_for_phase)
    all_chunks = []
    for phase in phases_chk:
        phase_items = [it for it in checklist if it.phase == phase]
        n_pages = max(1, -(-len(phase_items) // ITEMS_PER_SLIDE))  # ceiling division
        for pi in range(n_pages):
            chunk = phase_items[pi * ITEMS_PER_SLIDE : (pi + 1) * ITEMS_PER_SLIDE]
            all_chunks.append((chunk, phase, pi + 1, n_pages))

    for page_idx, (chunk, phase_label, phase_page, phase_total) in enumerate(all_chunks):
        sc = _add_branded_slide()
        slide_num = 5 + page_idx

        _add_rect(sc, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
                  _hex_to_rgb("FAFAF8"), _hex_to_rgb("FAFAF8"))

        # Phase color for header accent
        _PHASE_HDR_COLORS = {
            "Quality Structure":     "B91C1C",
            "Customer Requirements": "1565A0",
            "Workforce Actions":     "B45309",
            "Shutdown Prep":         "16283F",
            "Quality Readiness":     "16283F",
            "Quality Confirmation":  "1A5C1A",
        }
        phase_hex   = _PHASE_HDR_COLORS.get(phase_label, "16283F")
        phase_color = _hex_to_rgb(phase_hex)

        # Colored header bar in phase color
        _add_rect(sc, Inches(0.3), Inches(0.12), Inches(12.73), Inches(0.62),
                  phase_color, phase_color)

        # Slide title on colored bar
        suffix = f" ({phase_page}/{phase_total})" if phase_total > 1 else ""
        _add_text(sc, f"{phase_label.upper()}{suffix}  —  {part_name}",
                  Inches(0.45), Inches(0.18), Inches(12.3), Inches(0.5),
                  font_size=18, bold=True, color=WHITE)

        # Sub-labels
        _add_text(sc, prog_label.upper(),
                  Inches(0.45), Inches(0.76), Inches(4), Inches(0.22),
                  font_size=8, color=GRAY)
        _add_text(sc, f"{len(chunk)} action items",
                  Inches(9.8), Inches(0.76), Inches(3.2), Inches(0.22),
                  font_size=8, color=GRAY, align=PP_ALIGN.RIGHT)
        _add_rect(sc, Inches(0.3), Inches(0.98), Inches(12.73), Inches(0.02),
                  _hex_to_rgb("E5DFD3"), _hex_to_rgb("E5DFD3"))

        # Render items in 2 columns
        for i, item in enumerate(chunk):
            col = i % 2
            row = i // 2
            x   = Inches(0.30 + col * 6.38)
            y   = Inches(1.06 + row * 0.42)

            if y + Inches(0.42) > Inches(7.15):
                break

            # Row styling
            if item.done:
                bg, bdr, txt_color = _hex_to_rgb("F0FFF4"), GREEN, GREEN
            elif item.critical:
                bg, bdr, txt_color = _hex_to_rgb("FFF4F4"), RED, RED
            else:
                bg  = _hex_to_rgb("F9F6F2") if row % 2 == 0 else _hex_to_rgb("FFFFFF")
                bdr, txt_color = LIGHT_GRAY, DARK_GRAY

            _add_rect(sc, x, y, Inches(6.1), Inches(0.37), bg, bdr)

            # Phase color accent bar on left
            _add_rect(sc, x, y, Inches(0.05), Inches(0.37), phase_color, phase_color)

            # Status badge
            if item.done:   badge_bg, badge_char = GREEN, "✓"
            elif item.critical: badge_bg, badge_char = RED, "!"
            else:           badge_bg, badge_char = _hex_to_rgb("64748B"), "○"

            _add_rect(sc, x + Inches(0.1), y + Inches(0.1),
                      Inches(0.18), Inches(0.18), badge_bg, badge_bg)
            _add_text(sc, badge_char,
                      x + Inches(0.1), y + Inches(0.07), Inches(0.18), Inches(0.24),
                      font_size=6, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            # Step code
            _add_text(sc, item.step,
                      x + Inches(0.31), y + Inches(0.05), Inches(0.55), Inches(0.28),
                      font_size=7, bold=True, color=phase_color)

            # Item text — shorter if owner/date present
            has_meta = bool(getattr(item, "owner", "") or getattr(item, "due", ""))
            txt_w    = Inches(3.8) if has_meta else Inches(5.12)
            text = item.text[:60] + "…" if (has_meta and len(item.text) > 60) else (
                   item.text[:82] + "…" if len(item.text) > 82 else item.text)
            _add_text(sc, text,
                      x + Inches(0.9), y + Inches(0.05),
                      txt_w, Inches(0.3),
                      font_size=8.5, color=txt_color)

            # Owner — right side of row
            owner_val = getattr(item, "owner", "")
            due_val   = getattr(item, "due", "")
            if owner_val or due_val:
                meta_x = x + Inches(4.75)
                if owner_val:
                    _add_rect(sc, meta_x, y + Inches(0.05),
                              Inches(1.28), Inches(0.14),
                              _hex_to_rgb("F4EFE5"), _hex_to_rgb("D6CFB8"))
                    _add_text(sc, f"👤 {owner_val[:16]}",
                              meta_x + Inches(0.04), y + Inches(0.04),
                              Inches(1.22), Inches(0.16),
                              font_size=6.5, color=_hex_to_rgb("16283F"))
                if due_val:
                    _add_rect(sc, meta_x, y + Inches(0.21),
                              Inches(1.28), Inches(0.14),
                              _hex_to_rgb("F0FFF4"), _hex_to_rgb("86EFAC"))
                    # Format date nicely
                    try:
                        from datetime import datetime as _dt
                        due_fmt = _dt.fromisoformat(str(due_val)).strftime("%d %b %Y")
                    except Exception:
                        due_fmt = str(due_val)[:10]
                    _add_text(sc, f"📅 {due_fmt}",
                              meta_x + Inches(0.04), y + Inches(0.20),
                              Inches(1.22), Inches(0.16),
                              font_size=6.5, color=_hex_to_rgb("15803D"))

        if LOGO_COLOR.exists():
            sc.shapes.add_picture(str(LOGO_COLOR), Inches(10.3), Inches(7.02), Inches(2.6), Inches(0.53))
        _footer(sc, slide_num, part_name, prog_type)

    # ── Remove original template slides (keep only our generated slides) ───
    # The original template slides sit at indices 0.._original_slide_count-1.
    # Our slides were appended after them.
    # We remove them by deleting their <p:sldId> entries from the XML.
    if _original_slide_count > 0:
        try:
            from pptx.oxml.ns import qn as _qn
            sld_id_lst = prs.element.find(_qn("p:sldIdLst"))
            if sld_id_lst is not None:
                children = list(sld_id_lst)
                to_remove = min(_original_slide_count, len(children) - 1)
                # Never remove all children — keep at least our first generated slide
                for child in children[:to_remove]:
                    sld_id_lst.remove(child)
        except Exception as e:
            print(f"[report_generator] Template cleanup warning: {e}")

    # ── Return as bytes ───────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _build_kpis(prog_type: str, result: ScoreResult, meta: dict, customer_name: str) -> list:
    """
    Returns 6 KPI dicts with label/value/sub, tailored to program type.
    """
    primary_date = meta.get("primary_date", "—")
    sl_end_date  = meta.get("sl_end_date",  "—")

    if prog_type == "restart":
        return [
            {"label": "RISK SCORE",          "value": f"{result.score}/100",           "sub": "Calculated risk level"},
            {"label": "SHUTDOWN DURATION",   "value": meta.get("shutdown_duration","—"),"sub": "Shutdown length"},
            {"label": "SHUTDOWN TYPE",       "value": meta.get("shutdown_type","—"),    "sub": "Nature of shutdown"},
            {"label": "SHUTDOWN PREP",       "value": meta.get("prep_status","—"),      "sub": "22-action checklist"},
            {"label": "LPA COVERAGE",        "value": meta.get("lpa_coverage","—"),     "sub": "Per area before release"},
            {"label": "RESTART DATE",        "value": primary_date,                     "sub": "Production restart"},
        ]
    elif prog_type == "transfer":
        return [
            {"label": "RISK SCORE",       "value": f"{result.score}/100",       "sub": "BT risk assessment"},
            {"label": "PRA FORECAST",     "value": result.pra_forecast or "—",  "sub": f"Conformance ≈{result.conformance}%"},
            {"label": "EXECUTION STEPS",  "value": "31",                         "sub": "EAGP_4-4_MG_01-F01"},
            {"label": "BMG REQUIRED",     "value": "YES" if meta.get("customer")=="vw" else "N/A", "sub": "VW Build Sample Approval"},
            {"label": "PPAP LEVEL",       "value": result.ppap,                  "sub": "Customer PPAP"},
            {"label": "TARGET SOP",       "value": primary_date,                 "sub": "Transfer completion"},
        ]
    elif prog_type == "absence":
        gap = meta.get("hc_gap_pct", 0)
        return [
            {"label": "RISK SCORE",      "value": f"{result.score}/100",              "sub": "Workforce risk level"},
            {"label": "HEADCOUNT GAP",   "value": f"{gap:.0f}%",                      "sub": f"{meta.get('hc_affected',0)} positions"},
            {"label": "TURNOVER RATE",   "value": f"{meta.get('turnover_pct',0):.0f}%","sub": "Annualised"},
            {"label": "BACKUP NEEDED",   "value": str(meta.get("backup_count",0)),    "sub": "120% buffer recommendation"},
            {"label": "VOLUME IMPACTED", "value": str(meta.get("volume","—")),        "sub": "pcs/day at risk"},
            {"label": "EFFECTIVE DATE",  "value": primary_date,                       "sub": "Impact start date"},
        ]
    elif prog_type == "capacity":
        return [
            {"label": "RISK SCORE",       "value": f"{result.score}/100",           "sub": "Capacity change risk"},
            {"label": "DEMAND CHANGE",    "value": meta.get("delta_pct_str","—"),   "sub": "Current vs target"},
            {"label": "CURRENT VOLUME",   "value": str(meta.get("volume","—")),     "sub": "pcs/day"},
            {"label": "TARGET VOLUME",    "value": str(meta.get("volume_new","—")), "sub": "pcs/day"},
            {"label": "PPAP LEVEL",       "value": result.ppap,                     "sub": "Quality submission"},
            {"label": "CHANGE DATE",      "value": primary_date,                    "sub": "Activation date"},
        ]
    elif prog_type == "my_change":
        return [
            {"label": "RISK SCORE",     "value": f"{result.score}/100",           "sub": "Change risk level"},
            {"label": "ECCL STAGE",     "value": meta.get("eccl_stage","—"),      "sub": "EAEP Step 7–17"},
            {"label": "E&O COSTS",      "value": meta.get("eo_costs_label","—"),  "sub": "Disposal threshold"},
            {"label": "PPAP LEVEL",     "value": result.ppap,                     "sub": "Customer submission"},
            {"label": "CUSTOMER",       "value": customer_name,                   "sub": "OEM"},
            {"label": "IMPL. DATE",     "value": primary_date,                    "sub": "Change implementation"},
        ]
    elif prog_type == "new_plant":
        gp12  = meta.get("gp12_end", "TBD")
        ptype = meta.get("plant_type", "—").capitalize()
        return [
            {"label": "RISK SCORE",         "value": f"{result.score}/100",    "sub": "New plant risk level"},
            {"label": "MONITORING",         "value": "Continuous",             "sub": "No fixed SL window"},
            {"label": "INSPECTION MODE",    "value": result.inspection,        "sub": "GP12 / containment"},
            {"label": "PPAP LEVEL",         "value": result.ppap,              "sub": "Customer submission"},
            {"label": "PLANT TYPE",         "value": ptype,                    "sub": "Site classification"},
            {"label": "GP12 END",           "value": gp12,                     "sub": "GP12 exit date"},
        ]
    else:   # new program + generic
        return [
            {"label": "RISK SCORE",       "value": f"{result.score}/100",  "sub": "Program risk level"},
            {"label": "SAFE LAUNCH",      "value": f"{result.duration} days","sub": "Post-SOP window"},
            {"label": "PROJECTED FPY",    "value": result.fpy,              "sub": "First Pass Yield"},
            {"label": "INSPECTION MODE",  "value": result.inspection,       "sub": "Quality control method"},
            {"label": "PPAP LEVEL",       "value": result.ppap,             "sub": "Customer submission"},
            {"label": "SL WINDOW END",    "value": sl_end_date,             "sub": "Safe launch exit date"},
        ]


# ═════════════════════════════════════════════════════════════════════════════
# EXCEL GENERATOR
# ═════════════════════════════════════════════════════════════════════════════

def generate_excel(
    prog_type:  str,
    part_name:  str,
    result:     ScoreResult,
    checklist:  List[ChecklistItem],
    meta:       dict,
    customer_name: str = "Other",
) -> bytes:
    """
    Generates an Excel report and returns bytes.
    Sheet 1: Summary / Risk Score
    Sheet 2: Risk Factors
    Sheet 3: Checklist
    """
    wb = openpyxl.Workbook()   # creates a new workbook
    # Make the brand body font the workbook default so even unstyled cells
    # render in Barlow rather than openpyxl's Calibri fallback.
    wb._named_styles["Normal"].font = Font(name=BODY_FONT, size=11)

    # ── Styling helpers ──────────────────────────────────────────────────────
    # PatternFill fills a cell with a background color
    # "solid" = single color fill; fgColor = foreground (the color you see)
    def fill(hex_color: str):
        return PatternFill("solid", fgColor=hex_color)

    def hdr_font(size=11, bold=True, color="FFFFFF"):
        return Font(name=BODY_FONT, size=size, bold=bold, color=color)

    def body_font(size=10, bold=False, color="1F2937"):
        return Font(name=BODY_FONT, size=size, bold=bold, color=color)

    thin = Side(style="thin", color="D6CFB8")   # Side defines a border line
    box_border = Border(left=thin, right=thin, top=thin, bottom=thin)

    prog_label = PROG_LABELS.get(prog_type, "Safe Launch")

    # ── Sheet 1: Summary ─────────────────────────────────────────────────────
    ws1 = wb.active               # wb.active = the first sheet by default
    ws1.title = "Summary"

    ws1.column_dimensions["A"].width = 30   # set column width in characters
    ws1.column_dimensions["B"].width = 40

    headers = [
        ("Versigent Safe Launch Report", None, "16283F", 16),
        (f"Program Type: {prog_label}", None, "CD7925", 12),
        (f"Part / Project: {part_name}", None, "1F2937", 11),
        (f"Customer: {customer_name}", None, "1F2937", 11),
        (f"Risk Score: {result.score}/100", None, "1F2937", 11),
        (f"Risk Level: {result.risk}", None,
         "B91C1C" if result.risk=="HIGH" else "B45309" if result.risk=="MEDIUM" else "15803D", 12),
        ("Monitoring: Continuous — no fixed safe launch window", None, "1F2937", 11) \
            if prog_type == "new_plant" else \
            (f"Safe Launch Duration: {result.duration} days", None, "1F2937", 11),
        (f"FPY Target: {result.fpy}", None, "1F2937", 11),
        (f"Inspection Method: {result.inspection}", None, "1F2937", 11),
        (f"PPAP Level: {result.ppap}", None, "1F2937", 11),
        ("", None, "FFFFFF", 10),
        ("RECOMMENDATION", None, "CD7925", 11),
        (result.recommendation, None, "374151", 10),
    ]

    for row_idx, (text, _, color, size) in enumerate(headers, start=1):
        cell = ws1.cell(row=row_idx, column=1, value=text)
        cell.font = Font(name=BODY_FONT, size=size, bold=(row_idx <= 2), color=color)
        cell.alignment = Alignment(wrap_text=True)

    # ── Sheet 2: Risk Factors ─────────────────────────────────────────────────
    ws2 = wb.create_sheet("Risk Factors")
    ws2.column_dimensions["A"].width = 35
    ws2.column_dimensions["B"].width = 10
    ws2.column_dimensions["C"].width = 10
    ws2.column_dimensions["D"].width = 10

    # Header row
    for col, hdr in enumerate(["Risk Factor", "Score", "Max", "%"], start=1):
        c = ws2.cell(row=1, column=col, value=hdr)
        c.fill   = fill("16283F")
        c.font   = hdr_font()
        c.border = box_border
        c.alignment = Alignment(horizontal="center")

    # Data rows
    for row_idx, factor in enumerate(result.factors, start=2):
        # Alternate row shading — row_idx % 2 is 0 for even, 1 for odd
        row_fill = fill("F4EFE5") if row_idx % 2 == 0 else fill("FFFFFF")
        data = [factor.name, factor.value, factor.max, f"{factor.percent}%"]
        for col, val in enumerate(data, start=1):
            c = ws2.cell(row=row_idx, column=col, value=val)
            c.fill   = row_fill
            c.font   = body_font()
            c.border = box_border
            c.alignment = Alignment(horizontal="center" if col > 1 else "left")

    # ── Sheet 3: Checklist ────────────────────────────────────────────────────
    ws3 = wb.create_sheet("Checklist")
    ws3.column_dimensions["A"].width = 8
    ws3.column_dimensions["B"].width = 20
    ws3.column_dimensions["C"].width = 60
    ws3.column_dimensions["D"].width = 15

    for col, hdr in enumerate(["Step", "Phase", "Description", "Status"], start=1):
        c = ws3.cell(row=1, column=col, value=hdr)
        c.fill   = fill("CD7925")
        c.font   = hdr_font()
        c.border = box_border
        c.alignment = Alignment(horizontal="center")

    for row_idx, item in enumerate(checklist, start=2):
        status = "✓ Done" if item.done else ("⚠ Critical" if item.warn else "Pending")
        row_fill = (fill("F0FFF4") if item.done else
                    fill("FFF0EE") if item.warn else
                    fill("FAFAFA") if row_idx % 2 == 0 else fill("FFFFFF"))
        data = [item.step, item.phase, item.text, status]
        for col, val in enumerate(data, start=1):
            c = ws3.cell(row=row_idx, column=col, value=val)
            c.fill   = row_fill
            c.font   = body_font(color="15803D" if item.done else
                                        "B91C1C" if item.warn else "1F2937")
            c.border = box_border
            c.alignment = Alignment(wrap_text=True, vertical="top")

    # ── Sheet 4: Transposed Checklist ───────────────────────────────────────
    ws4 = wb.create_sheet("Transposed Plan")
    ws4.sheet_view.showGridLines = False

    phases_t  = get_phases(checklist)
    PHASE_HEX = {
        "Quality Structure":       "B91C1C",
        "Infrastructure Readiness":"5C3317",
        "Customer Requirements":   "1565A0",
        "Workforce Actions":       "B45309",
        "Shutdown Prep":           "16283F",
        "Quality Readiness":       "16283F",
        "Quality Confirmation":    "1A5C1A",
    }

    # Title row
    ws4.row_dimensions[1].height = 28
    n_phase_cols = len(phases_t)
    ws4.merge_cells(f"A1:{get_column_letter(n_phase_cols + 1)}1")
    tc = ws4["A1"]
    tc.value     = f"{prog_label.upper()}  |  {part_name}  |  {customer_name}"
    tc.font      = Font(name=BODY_FONT, size=13, bold=True, color="FFFFFF")
    tc.fill      = fill("16283F")
    tc.alignment = Alignment(horizontal="left", vertical="center")

    ws4.row_dimensions[2].height = 8   # spacer
    ws4.row_dimensions[3].height = 30  # phase headers

    ws4.column_dimensions["A"].width = 5  # row number col

    # Phase header row (row 3) — one column per phase
    for col_idx, phase in enumerate(phases_t):
        col_letter = get_column_letter(col_idx + 2)
        n_items_ph = sum(1 for i in checklist if i.phase == phase)
        ws4.column_dimensions[col_letter].width = max(30, min(52, n_items_ph * 2))

        hc = ws4[f"{col_letter}3"]
        hc.value     = f"{phase.upper()}  ({n_items_ph})"
        hc.font      = Font(name=BODY_FONT, size=10, bold=True, color="FFFFFF")
        hc.fill      = fill(PHASE_HEX.get(phase, "16283F"))
        hc.alignment = Alignment(horizontal="center", vertical="center")
        hc.border    = Border(
            left  =Side(style="medium", color="FFFFFF"),
            right =Side(style="medium", color="FFFFFF"),
            bottom=Side(style="thin",   color="D6CFB8"),
        )

    rn_hdr = ws4["A3"]
    rn_hdr.value = "#"; rn_hdr.font = Font(name=BODY_FONT, size=9, bold=True, color="FFFFFF")
    rn_hdr.fill  = fill("16283F"); rn_hdr.alignment = Alignment(horizontal="center", vertical="center")

    # Data rows — one row per position, one col per phase
    max_items_t = max((sum(1 for i in checklist if i.phase == ph) for ph in phases_t), default=0)
    for row_off in range(max_items_t):
        data_row = 4 + row_off
        ws4.row_dimensions[data_row].height = 52

        rn_c = ws4.cell(row=data_row, column=1, value=row_off + 1)
        rn_c.font      = Font(name=BODY_FONT, size=8, color="64748B")
        rn_c.fill      = fill("FAFAF8")
        rn_c.alignment = Alignment(horizontal="center", vertical="top")
        rn_c.border    = Border(right=Side(style="thin", color="E5DFD3"))

        for col_idx, phase in enumerate(phases_t):
            col_letter   = get_column_letter(col_idx + 2)
            phase_items_t = [i for i in checklist if i.phase == phase]
            cell = ws4.cell(row=data_row, column=col_idx + 2)
            ph_hex = PHASE_HEX.get(phase, "16283F")
            bg     = "F9F6F2" if row_off % 2 == 0 else "FFFFFF"

            if row_off < len(phase_items_t):
                item    = phase_items_t[row_off]
                owner_v = getattr(item, "owner", "")
                due_v   = getattr(item, "due",   "")
                lines   = [f"[{item.step}]  {item.text}"]
                if owner_v: lines.append(f"👤 {owner_v}")
                if due_v:
                    try:
                        from datetime import datetime as _dt
                        lines.append(f"📅 {_dt.fromisoformat(str(due_v)).strftime('%d %b %Y')}")
                    except Exception:
                        lines.append(f"📅 {str(due_v)[:10]}")

                cell.value = "\n".join(lines)
                if item.critical:
                    cell.font = Font(name=BODY_FONT, size=9, bold=True, color="B91C1C")
                    cell.fill = fill("FFF4F4")
                elif getattr(item, "done", False):
                    cell.font = Font(name=BODY_FONT, size=9, color="15803D")
                    cell.fill = fill("F0FFF4")
                else:
                    cell.font = Font(name=BODY_FONT, size=9, color="374151")
                    cell.fill = fill(bg)
                cell.border = Border(
                    left  =Side(style="medium", color=ph_hex),
                    right =Side(style="thin",   color="E5DFD3"),
                    top   =Side(style="thin",   color="E5DFD3"),
                    bottom=Side(style="thin",   color="E5DFD3"),
                )
            else:
                cell.fill   = fill("F4F4F4")
                cell.border = Border(
                    left  =Side(style="thin", color="E5DFD3"),
                    right =Side(style="thin", color="E5DFD3"),
                    top   =Side(style="thin", color="F4F4F4"),
                    bottom=Side(style="thin", color="F4F4F4"),
                )
            cell.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True)

    ws4.freeze_panes = "B4"

    # ── Sheet 5: Gantt Chart (same axes as Transposed) ────────────────────
    from datetime import datetime as _dt, date as _date, timedelta as _td

    def parse_due(d):
        try: return _dt.fromisoformat(str(d)).date()
        except Exception: return None

    dated_items = [(i, parse_due(getattr(i,"due","")))
                   for i in checklist if parse_due(getattr(i,"due",""))]

    if dated_items:
        ws5 = wb.create_sheet("Gantt Chart")
        ws5.sheet_view.showGridLines = False

        all_dates   = [d for _, d in dated_items]
        date_min    = min(all_dates)
        date_max    = max(all_dates)
        sop_raw     = meta.get("primary_date","")
        sop_date    = parse_due(sop_raw) if sop_raw and sop_raw not in ("—","") else None
        range_start = min(sop_date, date_min - _td(days=7)) if sop_date else date_min - _td(days=7)
        range_end   = date_max + _td(days=14)

        # Build weekly columns (Monday-anchored)
        weeks = []
        cur = range_start - _td(days=range_start.weekday())
        while cur <= range_end:
            weeks.append(cur)
            cur += _td(days=7)

        # Fixed columns: A=#  B=Phase  C=Item  D=Owner  E=Due
        FIXED = 5
        for ci, w in enumerate([4, 22, 44, 16, 12]):
            ws5.column_dimensions[get_column_letter(ci+1)].width = w
        for wi in range(len(weeks)):
            ws5.column_dimensions[get_column_letter(FIXED + 1 + wi)].width = 6

        total_cols = FIXED + len(weeks)

        # Row 1: Title
        ws5.row_dimensions[1].height = 28
        ws5.merge_cells(f"A1:{get_column_letter(total_cols)}1")
        gt = ws5["A1"]
        gt.value     = f"{prog_label.upper()} — GANTT CHART  |  {part_name}  |  {customer_name}"
        gt.font      = Font(name=BODY_FONT, size=13, bold=True, color="FFFFFF")
        gt.fill      = fill("16283F")
        gt.alignment = Alignment(horizontal="left", vertical="center")

        # Row 2: Month labels (merged per month)
        ws5.row_dimensions[2].height = 14
        month_groups = {}
        for wi, wk in enumerate(weeks):
            mk = wk.strftime("%B %Y")
            month_groups.setdefault(mk, []).append(wi)
        for mk, indices in month_groups.items():
            c_s = get_column_letter(FIXED + 1 + indices[0])
            c_e = get_column_letter(FIXED + 1 + indices[-1])
            if c_s != c_e:
                try: ws5.merge_cells(f"{c_s}2:{c_e}2")
                except Exception: pass
            mc = ws5[f"{c_s}2"]
            mc.value = mk
            mc.font  = Font(name=BODY_FONT, size=8, bold=True, color="FFFFFF")
            mc.fill  = fill("1F3553")
            mc.alignment = Alignment(horizontal="center", vertical="center")

        # Row 3: Column headers + week labels
        ws5.row_dimensions[3].height = 22
        today = _date.today()
        for ci, lbl in enumerate(["#","Phase","Action Item","Owner","Due"]):
            c = ws5.cell(row=3, column=ci+1, value=lbl)
            c.font      = Font(name=BODY_FONT, size=9, bold=True, color="FFFFFF")
            c.fill      = fill("CD7924")
            c.alignment = Alignment(horizontal="center", vertical="center")
            c.border    = Border(bottom=Side(style="thin", color="D6CFB8"),
                                 right=Side(style="thin", color="FFFFFF"))

        for wi, wk in enumerate(weeks):
            is_sop     = sop_date and wk <= sop_date < wk + _td(days=7)
            is_current = wk <= today < wk + _td(days=7)
            wk_fill = "F97316" if is_sop else ("BFDBFE" if is_current else "F4EFE5")
            wk_fc   = "FFFFFF" if is_sop else "374151"
            wc = ws5.cell(row=3, column=FIXED + 1 + wi, value=wk.strftime("%d/%m"))
            wc.font      = Font(name=BODY_FONT, size=7, bold=is_sop, color=wk_fc)
            wc.fill      = fill(wk_fill)
            wc.alignment = Alignment(horizontal="center", vertical="center")
            wc.border    = Border(
                left  =Side(style="thin", color="FFFFFF"),
                right =Side(style="thin", color="FFFFFF"),
                bottom=Side(style="thin", color="D6CFB8"),
            )

        # Row 4: SOP marker
        ws5.row_dimensions[4].height = 12
        if sop_date:
            sop_wi = next((wi for wi, wk in enumerate(weeks)
                           if wk <= sop_date < wk + _td(days=7)), None)
            if sop_wi is not None:
                sc2 = ws5.cell(row=4, column=FIXED + 1 + sop_wi, value="▼ SOP")
                sc2.font = Font(name=BODY_FONT, size=7, bold=True, color="B91C1C")
                sc2.fill = fill("FFE4E1")
                sc2.alignment = Alignment(horizontal="center", vertical="center")

        # Data rows — same order as Transposed sheet: sorted by phase then due date
        sorted_for_gantt = sorted(
            dated_items,
            key=lambda x: (
                list(PHASE_HEX.keys()).index(x[0].phase) if x[0].phase in PHASE_HEX else 99,
                x[1],
            )
        )

        # Phase section headers — insert a header row whenever phase changes
        output_rows = []  # list of ("header", phase) or ("item", item, due_date)
        prev_phase = None
        for item, due_d in sorted_for_gantt:
            if item.phase != prev_phase:
                output_rows.append(("header", item.phase))
                prev_phase = item.phase
            output_rows.append(("item", item, due_d))

        data_start_row = 5
        for row_off, entry in enumerate(output_rows):
            data_row = data_start_row + row_off
            ws5.row_dimensions[data_row].height = 22 if entry[0] == "header" else 26

            if entry[0] == "header":
                phase = entry[1]
                ph_hex = PHASE_HEX.get(phase, "16283F")
                ws5.merge_cells(f"A{data_row}:E{data_row}")
                pc = ws5[f"A{data_row}"]
                pc.value     = f"  {phase.upper()}"
                pc.font      = Font(name=BODY_FONT, size=9, bold=True, color="FFFFFF")
                pc.fill      = fill(ph_hex)
                pc.alignment = Alignment(horizontal="left", vertical="center")
                # Grey out gantt area for section header
                for wi in range(len(weeks)):
                    gc = ws5.cell(row=data_row, column=FIXED + 1 + wi)
                    gc.fill = fill(ph_hex)
                continue

            item, due_d = entry[1], entry[2]
            ph_hex  = PHASE_HEX.get(item.phase, "16283F")
            is_crit = item.critical
            owner_v = getattr(item, "owner", "")
            row_bg  = "FFF4F4" if is_crit else ("F9F6F2" if row_off % 2 == 0 else "FFFFFF")

            # Fixed columns
            for ci, (val, ha, bld) in enumerate([
                (row_off + 1,                                                      "center", False),
                (item.phase,                                                        "left",   False),
                (f"[{item.step}]  {item.text[:50]}{'…' if len(item.text)>50 else ''}", "left", is_crit),
                (owner_v,                                                           "center", False),
                (due_d.strftime("%d %b %Y") if due_d else "—",                    "center", True),
            ]):
                c = ws5.cell(row=data_row, column=ci + 1, value=val)
                c.font      = Font(name=BODY_FONT, size=8, bold=bld,
                                   color="B91C1C" if (is_crit and ci==2) else
                                          ph_hex   if ci in (1,4) else "374151")
                c.fill      = fill(row_bg)
                c.alignment = Alignment(horizontal=ha, vertical="center", wrap_text=(ha=="left"))
                c.border    = Border(
                    left  =Side(style="medium" if ci==0 else "thin",
                                color=ph_hex if ci==0 else "E5DFD3"),
                    right =Side(style="thin",   color="E5DFD3"),
                    top   =Side(style="thin",   color="E5DFD3"),
                    bottom=Side(style="thin",   color="E5DFD3"),
                )

            # Gantt bar — fill up to due week, milestone diamond on due week
            if due_d:
                due_wi = next((wi for wi, wk in enumerate(weeks)
                               if wk <= due_d < wk + _td(days=7)), None)
                sop_wi = next((wi for wi, wk in enumerate(weeks)
                               if sop_date and wk <= sop_date < wk + _td(days=7)), None)

                for wi in range(len(weeks)):
                    gc = ws5.cell(row=data_row, column=FIXED + 1 + wi)
                    gc.alignment = Alignment(horizontal="center", vertical="center")

                    # Determine fill
                    is_sop_col = (sop_wi is not None and wi == sop_wi)
                    is_due_col = (due_wi is not None and wi == due_wi)
                    bar_color  = "B91C1C" if is_crit else ph_hex

                    if is_due_col:
                        gc.fill  = fill(bar_color)
                        gc.value = "◆"
                        gc.font  = Font(name=BODY_FONT, size=9, bold=True, color="FFFFFF")
                    elif is_sop_col:
                        gc.fill  = fill("FFE4E1")
                        gc.value = ""
                    else:
                        gc.fill  = fill("F5F5F5" if row_off % 2 == 0 else "FAFAFA")
                        gc.value = ""

                    gc.border = Border(
                        left  =Side(style="thin", color="EEEEEE"),
                        right =Side(style="thin", color="EEEEEE"),
                        top   =Side(style="thin", color="EEEEEE"),
                        bottom=Side(style="thin", color="EEEEEE"),
                    )

        ws5.freeze_panes = "F5"

    # ── Return as bytes ───────────────────────────────────────────────────
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()
