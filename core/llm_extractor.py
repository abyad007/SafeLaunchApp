# core/llm_extractor.py
# ─────────────────────────────────────────────────────────────────────────────
# LLM-based procedure extraction engine.
# Compatible with vLLM and Ollama — both expose the same OpenAI-compatible API.
#
# Architecture:
#   1. check_llm_available()  — ping the server, return True/False
#   2. extract_with_llm()     — send procedure text, parse structured response
#   3. Automatic fallback     — if LLM fails for any reason, returns None
#                               so procedure_reader falls back to rule-based
#
# TISAX compliance:
#   - Zero data leaves the network — calls go to localhost or internal server only
#   - No API keys — local models don't require authentication
#   - All calls use standard HTTP (requests library, already in Python stdlib-ish)
#   - Timeout enforced — app never hangs waiting for LLM
#   - Graceful fallback — LLM failure never breaks the app
# ─────────────────────────────────────────────────────────────────────────────

from __future__ import annotations

import json
import re
import time
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple

# config is the single source of truth for all settings
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))
import config as cfg

# requests is included in Python standard environments and already used by Streamlit
try:
    import requests as _requests
    _REQUESTS_OK = True
except ImportError:
    _REQUESTS_OK = False


# ─────────────────────────────────────────────────────────────────────────────
# SYSTEM PROMPT
# This is the core of the LLM extraction quality.
# It is explicit, structured, and leaves no room for interpretation.
# ─────────────────────────────────────────────────────────────────────────────

_SYSTEM_PROMPT = """You are a quality engineering assistant specialised in automotive supplier quality procedures (IATF 16949, VDA, APQP).

Your task: extract ONLY concrete quality action items from the procedure text provided.

STRICT RULES — follow exactly:
1. Return ONLY a valid JSON array. No preamble, no explanation, no markdown.
2. Each element is an object with exactly these fields:
   - "step": string — original step number if present, else sequential "1", "2"...
   - "text": string — the action item, cleaned and complete (20–200 chars)
   - "phase": string — the section/slide title this item belongs to (empty string if unknown)
   - "critical": boolean — true if item involves PFMEA, PPAP, MSA, containment, SC, 100%, EOL, LPA, PRA, firewall, control plan, process release, customer approval
3. INCLUDE only items that are direct quality actions — things a quality engineer MUST DO.
4. EXCLUDE:
   - Role labels: R, S, CS, Mfg, Responsible, Accountable, Support, Inform
   - Status cells: YES, NO, N/A, OK, Done, Open
   - Table headers and column labels
   - Week/date references: Wk 20/26, Q3 2026
   - Conditions and questions: "if the process...", "is the part..."
   - Document references: "A. Quality Assurance Matrix – EAGP_4-3_CS_11"
   - Glossary definitions: "PFMEA – Process Failure Mode and Effects Analysis"
   - KPI metric labels: "% Completion of MQ1", "FTQ target"
   - Placeholder text: XXX, TBD, N/A
   - Footnotes and legend text
5. Do NOT invent or paraphrase — use the exact wording from the procedure.
6. Do NOT duplicate — if the same action appears multiple times, include it once.

Output format (JSON array only):
[
  {"step": "1", "text": "Review Customer and Tier 1 Expectations and Requirements", "phase": "Concept Phase", "critical": false},
  {"step": "2", "text": "Complete PFMEA for all critical characteristics before SOP", "phase": "Quality Readiness", "critical": true}
]"""


_USER_PROMPT_TEMPLATE = """Extract all quality action items from this procedure document.

PROCEDURE TEXT:
{procedure_text}

Return ONLY the JSON array. No other text."""


# ─────────────────────────────────────────────────────────────────────────────
# CONNECTION CHECK
# ─────────────────────────────────────────────────────────────────────────────

def check_llm_available() -> Tuple[bool, str]:
    """
    Pings the configured LLM server to check if it's reachable.

    Returns (available: bool, message: str)
    Fast timeout (3s) so the app doesn't hang on startup.

    For vLLM:  GET /v1/models
    For Ollama: GET /api/tags
    """
    if not cfg.LLM_ENABLED:
        return False, "LLM disabled in config.py (LLM_ENABLED = False)"

    if not _REQUESTS_OK:
        return False, "requests library not available"

    endpoint = cfg.LLM_ENDPOINT.rstrip("/")
    check_url = (
        f"{endpoint}/v1/models"     if cfg.LLM_PROVIDER == "vllm"
        else f"{endpoint}/api/tags" if cfg.LLM_PROVIDER == "ollama"
        else f"{endpoint}/v1/models"
    )

    try:
        resp = _requests.get(check_url, timeout=3)
        if resp.status_code == 200:
            return True, f"LLM server reachable at {endpoint}"
        else:
            return False, f"LLM server returned HTTP {resp.status_code}"
    except _requests.exceptions.ConnectionError:
        return False, f"LLM server not reachable at {endpoint} — using rule-based extraction"
    except _requests.exceptions.Timeout:
        return False, f"LLM server timeout at {endpoint}"
    except Exception as e:
        return False, f"LLM check error: {str(e)[:80]}"


# ─────────────────────────────────────────────────────────────────────────────
# CORE EXTRACTION
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class LLMItem:
    """One item returned by the LLM."""
    step:     str
    text:     str
    phase:    str  = ""
    critical: bool = False


def _build_request_payload(procedure_text: str) -> dict:
    """
    Builds the API request body.
    Uses the OpenAI-compatible /v1/chat/completions format.
    Both vLLM and recent Ollama versions support this endpoint.
    """
    return {
        "model":       cfg.LLM_MODEL,
        "messages": [
            {"role": "system", "content": _SYSTEM_PROMPT},
            {"role": "user",   "content": _USER_PROMPT_TEMPLATE.format(
                procedure_text=procedure_text[:cfg.LLM_MAX_INPUT_CHARS]
            )},
        ],
        "temperature":  cfg.LLM_TEMPERATURE,
        "max_tokens":   cfg.LLM_MAX_TOKENS,
        "stream":       False,
    }


def _call_llm(payload: dict) -> Optional[str]:
    """
    Makes the HTTP POST to the LLM server.
    Returns raw response text or None on any failure.
    """
    endpoint = cfg.LLM_ENDPOINT.rstrip("/")

    # Both vLLM and Ollama (v0.3+) support /v1/chat/completions
    url = f"{endpoint}/v1/chat/completions"

    try:
        resp = _requests.post(
            url,
            json=payload,
            headers={"Content-Type": "application/json"},
            timeout=cfg.LLM_TIMEOUT,
        )
        resp.raise_for_status()
        data = resp.json()

        # Standard OpenAI-compatible response format
        return data["choices"][0]["message"]["content"]

    except _requests.exceptions.Timeout:
        print(f"[llm_extractor] Timeout after {cfg.LLM_TIMEOUT}s")
        return None
    except _requests.exceptions.ConnectionError:
        print(f"[llm_extractor] Cannot connect to {endpoint}")
        return None
    except KeyError as e:
        print(f"[llm_extractor] Unexpected response format: {e}")
        return None
    except Exception as e:
        print(f"[llm_extractor] Error: {e}")
        return None


def _parse_llm_response(raw: str) -> Optional[List[LLMItem]]:
    """
    Parses the LLM response into a list of LLMItems.

    The LLM is instructed to return pure JSON, but in practice it
    sometimes wraps it in markdown code fences or adds preamble.
    This parser handles all common deviation patterns gracefully.
    """
    if not raw or not raw.strip():
        return None

    # Strip markdown code fences if present: ```json ... ``` or ``` ... ```
    cleaned = re.sub(r"```(?:json)?\s*", "", raw).strip()
    cleaned = re.sub(r"```\s*$", "", cleaned).strip()

    # Find the JSON array — look for [ ... ] even if there's text before/after
    match = re.search(r"\[.*\]", cleaned, re.DOTALL)
    if not match:
        print("[llm_extractor] No JSON array found in response")
        return None

    json_str = match.group(0)

    try:
        data = json.loads(json_str)
    except json.JSONDecodeError as e:
        # Try to fix common LLM JSON errors: trailing commas, single quotes
        json_str_fixed = re.sub(r",\s*([}\]])", r"\1", json_str)   # trailing commas
        json_str_fixed = json_str_fixed.replace("'", '"')           # single → double quotes
        try:
            data = json.loads(json_str_fixed)
        except json.JSONDecodeError:
            print(f"[llm_extractor] JSON parse error: {e}")
            return None

    if not isinstance(data, list):
        print(f"[llm_extractor] Expected list, got {type(data)}")
        return None

    items: List[LLMItem] = []
    for i, entry in enumerate(data):
        if not isinstance(entry, dict):
            continue

        text = str(entry.get("text", "")).strip()
        if not text or len(text) < 10:
            continue

        items.append(LLMItem(
            step     = str(entry.get("step", str(i + 1))),
            text     = text,
            phase    = str(entry.get("phase", "")),
            critical = bool(entry.get("critical", False)),
        ))

    return items if items else None


def _validate_items(items: List[LLMItem]) -> List[LLMItem]:
    """
    Post-processing validation on LLM output.
    Removes any items that still look like noise despite the prompt.
    Uses the same rules as the rule-based extractor for consistency.
    """
    # Import here to avoid circular dependency
    from core.procedure_reader import _is_quality_action

    validated = []
    seen_keys: set = set()

    for item in items:
        # Deduplicate
        key = item.text[:55].lower().strip()
        if key in seen_keys:
            continue
        seen_keys.add(key)

        # Minimum length
        if len(item.text) < 15:
            continue

        # Run through the same linguistic filter as rule-based
        # This catches cases where the LLM ignores the prompt
        if not _is_quality_action(item.text):
            # LLM is usually right — only reject obvious non-actions
            # (very short, pure acronym, ALL CAPS header)
            if (len(item.text) < 20
                    or item.text.isupper()
                    or re.match(r"^[A-Z\s/]{1,30}$", item.text)):
                continue

        validated.append(item)

    return validated


# ─────────────────────────────────────────────────────────────────────────────
# TEXT EXTRACTION FROM FILES
# (reused from procedure_reader to avoid importing pptx/openpyxl twice)
# ─────────────────────────────────────────────────────────────────────────────

def _extract_full_text(filepath: Path) -> str:
    """
    Extracts ALL text from a procedure file for sending to the LLM.
    More permissive than the rule-based extractor — the LLM decides what to keep.
    Respects LLM_MAX_INPUT_CHARS to stay within model context window.
    """
    ext  = filepath.suffix.lower()
    text_parts: List[str] = []
    chars = 0

    try:
        if ext in {".pptx", ".pptm"}:
            from pptx import Presentation
            prs = Presentation(str(filepath))

            for slide in prs.slides:
                # Get slide title for context
                title = ""
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        try:
                            if shape.placeholder_format.idx == 0:
                                title = shape.text_frame.text.strip()
                                break
                        except Exception:
                            pass

                if title:
                    text_parts.append(f"\n## {title}\n")

                # Tables first (structured content)
                for shape in slide.shapes:
                    if shape.shape_type == 19:  # TABLE
                        for row in shape.table.rows:
                            row_texts = []
                            for cell in row.cells:
                                ct = cell.text.strip()
                                if ct:
                                    row_texts.append(ct)
                            if row_texts:
                                line = " | ".join(row_texts)
                                text_parts.append(line)
                                chars += len(line)

                # Text frames
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    if shape.is_placeholder:
                        try:
                            if shape.placeholder_format.idx == 0:
                                continue  # skip title — already added
                        except Exception:
                            pass
                    t = shape.text_frame.text.strip()
                    if t:
                        text_parts.append(t)
                        chars += len(t)

                if chars >= cfg.LLM_MAX_INPUT_CHARS:
                    break

        elif ext in {".xlsx", ".xlsm", ".xls"}:
            import openpyxl
            wb = openpyxl.load_workbook(str(filepath), read_only=False, data_only=True)

            for sname in wb.sheetnames:
                if any(skip in sname.lower()
                       for skip in ["template", "reference", "revision", "index", "legend"]):
                    continue

                text_parts.append(f"\n## {sname}\n")
                sheet = wb[sname]

                for row in sheet.iter_rows(min_row=1, values_only=True):
                    row_texts = [str(c).strip() for c in row
                                 if c is not None and str(c).strip()
                                 and str(c).strip() not in ("None", "")]
                    if row_texts:
                        line = " | ".join(row_texts)
                        text_parts.append(line)
                        chars += len(line)
                        if chars >= cfg.LLM_MAX_INPUT_CHARS:
                            break

                if chars >= cfg.LLM_MAX_INPUT_CHARS:
                    break

            wb.close()

    except Exception as e:
        print(f"[llm_extractor] Text extraction error: {e}")

    full_text = "\n".join(text_parts)
    return full_text[:cfg.LLM_MAX_INPUT_CHARS]


# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API
# ─────────────────────────────────────────────────────────────────────────────

def extract_with_llm(filepath: Path) -> Optional[List[LLMItem]]:
    """
    Main entry point for LLM-based extraction.

    Returns a list of LLMItems if successful, or None if:
    - LLM server is not reachable
    - LLM returns garbage / unparseable output
    - LLM returns fewer items than LLM_MIN_ITEMS_THRESHOLD

    When None is returned, the caller (procedure_reader) falls back
    to the rule-based extractor automatically.

    Usage:
        items = extract_with_llm(Path("data/procedures/HOGP.xlsx"))
        if items:
            # use LLM items
        else:
            # use rule-based fallback
    """
    if not cfg.LLM_ENABLED:
        return None

    if not _REQUESTS_OK:
        return None

    # Extract text from file
    procedure_text = _extract_full_text(filepath)
    if not procedure_text or len(procedure_text) < 100:
        print(f"[llm_extractor] Too little text extracted from {filepath.name}")
        return None

    # Build and send request
    t_start  = time.time()
    payload  = _build_request_payload(procedure_text)
    raw_resp = _call_llm(payload)
    t_elapsed = time.time() - t_start

    if raw_resp is None:
        return None

    # Parse response
    items = _parse_llm_response(raw_resp)
    if items is None:
        print(f"[llm_extractor] Could not parse LLM response for {filepath.name}")
        return None

    # Validate and clean
    items = _validate_items(items)

    # Safety check — if too few items, probably something went wrong
    if len(items) < cfg.LLM_MIN_ITEMS_THRESHOLD:
        print(f"[llm_extractor] Only {len(items)} items returned — falling back to rule-based")
        return None

    print(f"[llm_extractor] {filepath.name}: {len(items)} items in {t_elapsed:.1f}s")
    return items


def get_llm_status() -> dict:
    """
    Returns a status dict for display in the Streamlit sidebar.
    Called once on app load and cached.

    Returns:
        {
            "available": bool,
            "message":   str,
            "provider":  str,
            "model":     str,
            "endpoint":  str,
            "enabled":   bool,
        }
    """
    available, message = check_llm_available()
    return {
        "available": available,
        "message":   message,
        "provider":  cfg.LLM_PROVIDER,
        "model":     cfg.LLM_MODEL,
        "endpoint":  cfg.LLM_ENDPOINT,
        "enabled":   cfg.LLM_ENABLED,
    }
